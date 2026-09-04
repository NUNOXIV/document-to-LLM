#!/usr/bin/env python3
"""ACSOS document-to-LLM — Abweichungspruefung Quelle gegen Extrakt.

Vergleicht den Textlayer des Quell-PDFs (Docling-PDF-Backend, ohne ML-Modelle)
Wort fuer Wort mit dem erzeugten Markdown. Ergebnis ist eine Deckungsquote plus
die Liste der Woerter, die im Extrakt fehlen — seitengenau.

Damit ist belegbar, dass nichts verloren gegangen ist, statt es zu hoffen.

    python verify.py output/iso-27001.md --source input/ISO-27001.pdf
    python verify.py output/iso-27001.md --source input/ISO-27001.pdf --min-coverage 99.5
"""

from __future__ import annotations

import re
import sys
import unicodedata
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

import click

# Zeichen, die zwischen PDF-Textlayer und Markdown-Serialisierung typografisch
# variieren duerfen, ohne dass es eine inhaltliche Abweichung waere.
_QUOTES = {
    "‘": "'", "’": "'", "‚": "'", "‛": "'",
    "“": '"', "”": '"', "„": '"', "‟": '"',
    "‐": "-", "‑": "-", "‒": "-", "–": "-", "—": "-",
    "−": "-", "­": "", " ": " ",
}


def normalize(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    for src, dst in _QUOTES.items():
        text = text.replace(src, dst)
    return text.casefold()


def tokenize(text: str) -> list[str]:
    """Woerter und Zahlen; Satzzeichen und Layout-Artefakte fallen weg."""
    text = normalize(text)
    # Am Zeilenende getrennte Woerter zusammenfuehren (Silbentrennung im PDF).
    text = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", text)
    return re.findall(r"[0-9a-zA-ZÀ-ɏ]+(?:[.,][0-9]+)*", text)


def markdown_tokens(md_path: Path) -> list[str]:
    text = md_path.read_text(encoding="utf-8")
    text = re.sub(r"\A---\n.*?\n---\n", "", text, flags=re.S)   # Front-Matter
    text = re.sub(r"<!--.*?-->", " ", text, flags=re.S)          # Seitenmarker/Hinweise
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", " ", text)            # Bildplatzhalter
    text = re.sub(r"^\s*\|[\s:|-]+\|\s*$", " ", text, flags=re.M)  # Tabellentrennzeilen
    return tokenize(text)


def _mask(line: str) -> str:
    """Zeile fuer den Kopf-/Fusszeilen-Vergleich normalisieren (Ziffern maskiert,
    damit 'Seite 3' und 'Seite 4' als dieselbe laufende Zeile gelten)."""
    return re.sub(r"\d+", "#", normalize(line)).strip()


def pdf_pages(pdf_path: Path) -> tuple[dict[int, list[str]], dict[int, list[str]]]:
    """Textlayer des PDFs ueber das Docling-Backend (keine ML-Modelle noetig).

    Rueckgabe: (Inhaltstokens je Seite, Kopf-/Fusszeilen-Tokens je Seite).
    Laufende Kopf- und Fusszeilen werden von Docling bewusst entfernt; sie
    duerfen die Deckungsquote nicht druecken und werden separat ausgewiesen.
    """
    from docling.backend.docling_parse_backend import DoclingParseDocumentBackend
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.document import InputDocument
    from docling_core.types.doc.page import TextCellUnit

    # InputDocument instanziiert das Backend bereits selbst; ein zweites wuerde
    # dieselbe Datei ein zweites Mal offen halten.
    in_doc = InputDocument(
        path_or_stream=pdf_path, format=InputFormat.PDF,
        backend=DoclingParseDocumentBackend, filename=pdf_path.name,
    )
    backend = in_doc._backend
    lines_by_page: dict[int, list[str]] = {}
    try:
        n = backend.page_count()
        for i in range(n):
            seg = backend.load_page(i).get_segmented_page()
            lines_by_page[i + 1] = [c.text for c in seg.iterate_cells(TextCellUnit.LINE)]
    finally:
        backend.unload()

    # Laufende Kopf-/Fusszeilen: gleiche (ziffernmaskierte) Zeile auf vielen Seiten,
    # jeweils am oberen oder unteren Rand des Seiteninhalts.
    n = len(lines_by_page)
    seen: Counter[str] = Counter()
    for lines in lines_by_page.values():
        seen.update({_mask(z) for z in lines if _mask(z) and len(tokenize(z)) <= 12})
    threshold = max(2, int(0.6 * n))
    running = {m for m, c in seen.items() if c >= threshold} if n >= 2 else set()

    content: dict[int, list[str]] = {}
    boiler: dict[int, list[str]] = {}
    for page_no, lines in lines_by_page.items():
        keep, drop = [], []
        for line in lines:
            (drop if _mask(line) in running else keep).append(line)
        content[page_no] = tokenize(" ".join(keep))
        boiler[page_no] = tokenize(" ".join(drop))
    return content, boiler


# --------------------------------------------------------------------------
# Verlorene Bindestriche
# --------------------------------------------------------------------------
# Docling loest die Silbentrennung am Zeilenende auf, indem es den Trennstrich
# entfernt. Das ist richtig fuer ein Wort, das nur der Umbruch getrennt hat
# ("Informations-\nsicherheit"), und falsch fuer ein Wort, das den Bindestrich
# selbst traegt ("IKT-\nSysteme"). Aus "IKT-Systemen" wird dann "IKTSystemen".
#
# Die Deckungspruefung sah das nicht: tokenize() entfernt denselben Trennstrich
# auf der Quellseite, beide Seiten hiessen "iktsystemen", die Deckung blieb
# 100 %. Gefunden hat es erst der Abgleich gegen den amtlichen Gesetzestext.
#
# Repariert wird nur mit Beleg. Ein Wort wird angefasst, wenn seine Form ohne
# Bindestrich in der Quelle NICHT vorkommt und die Form mit Bindestrich dort
# vorkommt. Damit bleibt echtes Binnenmajuskel ("OpenLDAP", "PowerShell")
# unangetastet: es steht so in der Quelle.
_WORT = re.compile(r"[A-Za-zÄÖÜäöüß0-9]+")

# Kuerzere Woerter bleiben unangetastet: bei drei, vier Zeichen ist die Gefahr
# eines zufaelligen Belegtreffers zu gross.
MIN_WORTLAENGE = 6


# Zeichen, die im Textlayer stehen, aber keinen Text bedeuten: Nichtzeichen und
# das Ersatzzeichen. Sie entstehen, wenn die Schrift des PDF einen Codepunkt
# nicht abbildet — im BSIG-Druck trifft es den geschuetzten Bindestrich U+2011,
# der als U+FFFE ankommt. Docling wirft das Zeichen weg, und aus "IKT-Systemen"
# wird "IKTSystemen": ein Wort, das es nirgends gibt.
UNLESBAR = "\ufffe\uffff\ufffd"
_UNLESBAR_IM_WORT = re.compile(rf"(?<=\w)[{UNLESBAR}](?=\w)")


def quelltext(pdf_path: Path) -> str:
    """Roher Textlayer der Quelle (Docling-unabhaengig, ohne ML-Modelle)."""
    import pypdfium2

    doc = pypdfium2.PdfDocument(pdf_path)
    try:
        return "".join(doc[i].get_textpage().get_text_range() for i in range(len(doc)))
    finally:
        doc.close()


def unlesbar_im_wort(text: str) -> int:
    """Zahl der unlesbaren Zeichen, die zwischen zwei Wortzeichen stehen."""
    return len(_UNLESBAR_IM_WORT.findall(text))


def quelle_kompakt(pdf_path: Path | str, text: str | None = None) -> str:
    """Textlayer der Quelle ohne jeden Zwischenraum.

    Dient dem ersten Beleg: kommt ein Wort hier vor, ist es richtig, egal wie
    es aussieht.

    Unlesbare Zeichen werden NICHT zu Bindestrichen gemacht. Der erste Anlauf
    tat das — mit dem Argument, im BSIG-Druck seien alle 24 solchen Stellen
    laut amtlichem XML Bindestriche. Das stimmte fuer das BSIG und nur dafuer:
    in einer anderen Datei desselben Bestandes stehen 1639 solche Zeichen, und
    dort sind es Trennstriche am Zeilenende ("Ab-nahme"). Dasselbe Zeichen,
    zwei Bedeutungen, kein Merkmal, das sie unterscheidet. Also wird geraten
    nicht — sondern gemeldet (siehe unlesbar_im_wort).
    """
    roh = text if text is not None else quelltext(Path(pdf_path))
    return re.sub(r"\s+", "", roh)


def verlorene_bindestriche(body: str, kompakt: str, zusammenhaengend: str = "") -> dict[str, str]:
    """Woerter im Extrakt, denen ein Bindestrich der Quelle fehlt.

    Rueckgabe: {"IKTSystemen": "IKT-Systemen", ...} — jeder Eintrag doppelt
    belegt.

    Zwei Belege, und der zweite ist der entscheidende:

    1. Die Form OHNE Bindestrich kommt in der Quelle nirgends vor. Sonst waere
       sie richtig — auch Binnenmajuskel wie "OpenLDAP" faellt darunter.
    2. Die Form MIT Bindestrich steht in der Quelle ZUSAMMENHAENGEND, also
       ohne Zeilenumbruch dazwischen.

    Ohne Beleg 2 kehrt die Reparatur die Silbentrennung um: in einem Dokument
    stand "Ab-\nnahme", der Extrakt hatte richtig "Abnahme" daraus gemacht, und
    die erste Fassung dieser Funktion haette daraus wieder "Ab-nahme" gemacht —
    93 solcher Fehlalarme in einer einzigen Datei. Der Unterschied ist genau
    die Zusammenhaengendheit: ein Bindestrich, der zum Wort gehoert, steht
    irgendwo auch mitten in der Zeile; ein Trennstrich steht nur am Zeilenende.
    """
    zusammenhaengend = zusammenhaengend or kompakt
    treffer: dict[str, str] = {}
    for wort in set(_WORT.findall(body)):
        if len(wort) < MIN_WORTLAENGE or wort in kompakt:
            continue
        # Jede Trennstelle wird geprueft, nicht nur der Wechsel von klein zu
        # gross: "IKTSystemen" und "Beratungsoder" haben keinen solchen
        # Wechsel und waeren sonst unsichtbar.
        for i in range(1, len(wort)):
            mit = wort[:i] + "-" + wort[i:]
            if mit in zusammenhaengend:
                treffer[wort] = mit
                break
    return treffer


def zusammenhaengende_quelle(pdf_path: Path | str, text: str | None = None) -> str:
    """Quelltext unveraendert, mit erhaltenen Umbruechen.

    Der zweite Beleg beruht genau darauf: ein am Zeilenende getrenntes Wort ist
    hier NICHT zusammenhaengend zu finden, ein Wort mit eigenem Bindestrich
    schon. Deshalb wird hier nichts ersetzt — auch kein unlesbares Zeichen.
    """
    return text if text is not None else quelltext(Path(pdf_path))


def repariere_bindestriche(body: str, kompakt: str,
                           zusammenhaengend: str = "") -> tuple[str, dict[str, str]]:
    """Setzt belegte Bindestriche zurueck. Ohne Beleg bleibt alles, wie es ist."""
    treffer = verlorene_bindestriche(body, kompakt, zusammenhaengend)
    for falsch, richtig in treffer.items():
        body = re.sub(rf"\b{re.escape(falsch)}\b", richtig, body)
    return body, treffer


def normalize_ooxml_styles(src: Path, workdir: Path) -> Path | None:
    """Repariert leere <fill/>-Elemente in der styles.xml einer Office-Datei.

    Web-Exporte (z. B. der CIS Controls Navigator) schreiben `<fill />` ohne
    patternFill. Excel oeffnet das, openpyxl — und damit Doclings
    MsExcelDocumentBackend — bricht mit `expected Fill` ab. Repariert wird
    ausschliesslich das Stylesheet auf Container-Ebene; Zellinhalte, Formeln
    und Blattstruktur werden unveraendert uebernommen. Gibt den Pfad der
    reparierten Kopie zurueck, oder None, wenn nichts zu reparieren war.
    """
    import zipfile

    empty_fill = re.compile(r"<((?:\w+:)?)fill\s*/>")
    target = workdir / src.name
    replaced = 0
    try:
        with zipfile.ZipFile(src) as zin:
            names = zin.namelist()
            if not any(n.endswith("styles.xml") for n in names):
                return None
            with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename.endswith("styles.xml"):
                        text = data.decode("utf-8", "replace")
                        fixed, n = empty_fill.subn(
                            r"<\1fill><\1patternFill patternType='none'/></\1fill>", text)
                        replaced += n
                        data = fixed.encode("utf-8")
                    zout.writestr(item, data)
    except (OSError, zipfile.BadZipFile):
        return None
    if not replaced:
        target.unlink(missing_ok=True)
        return None
    return target


def readable_office_source(path: Path) -> tuple[Path, object]:
    """Liefert einen fuer openpyxl lesbaren Pfad zur Office-Quelle.

    Normalfall: die Quelle selbst. Bei defektem Stylesheet (leere <fill/>-
    Elemente aus Web-Exporten) eine stylesbereinigte Kopie in einem
    TemporaryDirectory — dessen Handle wird mit zurueckgegeben, damit der
    Aufrufer es offen haelt, solange gelesen wird.
    """
    import tempfile
    from openpyxl import load_workbook

    try:
        load_workbook(path, read_only=True, data_only=True).close()
        return path, None
    except TypeError:
        tmp = tempfile.TemporaryDirectory(prefix="acsos-ooxml-")
        repaired = normalize_ooxml_styles(path, Path(tmp.name))
        if repaired is None:
            tmp.cleanup()
            raise
        return repaired, tmp


def office_pages(path: Path) -> tuple[dict[int, list[str]], dict[int, list[str]]]:
    """Quelltext aus Office-Formaten — je Blatt, Abschnitt oder Folie.

    Gelesen wird mit dem Standardleser des jeweiligen Formats, also unabhaengig
    von Docling. Genau das ist der Sinn der Pruefung: die Extraktion macht
    Docling, der Abgleich kommt aus einer zweiten Quelle.
    """
    suffix = path.suffix.lower()
    pages: dict[int, list[str]] = {}

    if suffix in (".xlsx", ".xlsm"):
        from openpyxl import load_workbook

        readable, _tmp = readable_office_source(path)
        wb = load_workbook(readable, read_only=True, data_only=True)
        try:
            for i, ws in enumerate(wb.worksheets, 1):
                words: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            words.extend(tokenize(str(cell)))
                pages[i] = words
        finally:
            wb.close()

    elif suffix == ".docx":
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))
        words = []
        for para in doc.paragraphs:
            words.extend(tokenize(para.text))
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    words.extend(tokenize(cell.text))
        pages[1] = words

    elif suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            words = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    words.extend(tokenize(shape.text_frame.text))
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            words.extend(tokenize(cell.text))
            pages[i] = words

    return pages, {p: [] for p in pages}


def source_pages(path: Path) -> tuple[dict[int, list[str]], dict[int, list[str]]]:
    """Quelltext je Seite (PDF) bzw. je Blatt/Folie (Office)."""
    if path.suffix.lower() == ".pdf":
        return pdf_pages(path)
    return office_pages(path)


@dataclass
class VerifyResult:
    source: str
    extract: str
    source_tokens: int = 0
    extract_tokens: int = 0
    missing_tokens: int = 0
    # Kein Vorgabewert 100.0: ein ungeprueftes Ergebnis darf nicht wie ein
    # perfektes aussehen. Ist keine Pruefung moeglich, bleibt das Feld None,
    # und wer die Zahl liest, ohne 'note' zu beachten, bekommt keinen Wert
    # geschenkt, sondern faellt auf.
    coverage: float | None = None
    per_page: dict[int, float] = field(default_factory=dict)
    missing_sample: list[str] = field(default_factory=list)
    worst_pages: list[tuple[int, float]] = field(default_factory=list)
    boilerplate_tokens: int = 0
    rejoined_tokens: int = 0
    note: str | None = None

    @property
    def ok_for(self) -> float:
        return self.coverage


def verify(pdf_path: Path, md_path: Path) -> VerifyResult:
    res = VerifyResult(source=str(pdf_path), extract=str(md_path))
    pages, boiler = source_pages(pdf_path)
    src_all = [t for toks in pages.values() for t in toks]
    out = markdown_tokens(md_path)
    res.source_tokens = len(src_all)
    res.extract_tokens = len(out)
    res.boilerplate_tokens = sum(len(t) for t in boiler.values())

    # Eine Deckung ist nur so viel wert wie die Grundlage, gegen die sie
    # rechnet. Ist die Grundlage leer, faellt das auf. Gefaehrlich ist der Fall
    # dazwischen: ein gescanntes PDF, dessen Textlayer nur ein paar Streuzeichen
    # aus Kopf- oder Fusszeilen enthaelt. Dann rechnet der Vergleich gegen eine
    # Handvoll Woerter und meldet 100 % — ein perfekter Wert ueber eine
    # Stichprobe, die nichts aussagt. Beobachtet an einem 291-Seiten-Scan:
    # 2 Referenzwoerter gegen 195773 Extraktwoerter, Ergebnis "Deckung 100,0 %".
    # Deshalb gilt die Grundlage erst ab einem absoluten Mindestumfang und einem
    # plausiblen Verhaeltnis zum Extrakt als brauchbar.
    # Rein relativ gemessen, nicht an einer absoluten Untergrenze: ein kleines
    # Tabellenblatt hat legitim wenig Text, und sein Extrakt ist genauso klein
    # — das Verhaeltnis bleibt gesund. Verdaechtig ist allein, wenn der Extrakt
    # umfangreich ist und die Vergleichsgrundlage dazu verschwindet.
    # Nur fuer PDFs: das Risiko, gegen das diese Pruefung schuetzt, ist
    # OCR-Text ueber einem fast leeren Textlayer — und OCR gibt es nur hier.
    # Office-Formate haben immer einen echten Leser; liefert er ueberhaupt
    # Text, ist die Deckung ueber diesen Text aussagekraeftig.
    #
    # Das ist nicht theoretisch: A3_Modellierung_Recplast_GmbH.xlsx verbindet
    # Zellen ueber 45 Spalten, Docling schreibt den Wert in jede davon. Aus
    # 4995 Quellwoertern werden 114056 Extraktwoerter — ein Verhaeltnis von
    # 4,4 %, allein durch Wiederholung. Ohne diese Einschraenkung haette die
    # Pruefung dieser Datei ihre berechtigte Deckungszahl genommen.
    MIN_ANTEIL = 0.05
    EXTRAKT_ERHEBLICH = 200
    zu_duenn = bool(src_all) and pdf_path.suffix.lower() == ".pdf" \
        and len(out) > EXTRAKT_ERHEBLICH \
        and len(src_all) < MIN_ANTEIL * len(out)

    if not src_all or zu_duenn:
        if zu_duenn:
            res.note = (
                f"Textlayer zu duenn fuer einen Wortvergleich: {len(src_all)} "
                f"Referenzwoerter gegen {len(out)} Woerter im Extrakt. Eine "
                f"Deckungszahl waere hier eine Scheingenauigkeit — sie wuerde "
                f"gegen Streuzeichen aus Kopf- oder Fusszeilen rechnen, nicht "
                f"gegen den Inhalt. Behandelt wie ein Scan ohne Textlayer: der "
                f"Text stammt aus der Zeichenerkennung und ist nicht geprueft."
            )
        else:
            res.note = (
                "Kein Textlayer im PDF (gescannt) — ein Wortvergleich ist nicht "
                "moeglich. Extraktion mit --ocr on pruefen."
                if pdf_path.suffix.lower() == ".pdf" else
                f"Format {pdf_path.suffix} wird fuer den Wortvergleich nicht "
                f"unterstuetzt — Extrakt nicht gegen die Quelle geprueft."
            )
        return res

    stream = "".join(out)  # fuer Woerter, die im PDF ueber einen Zeilenumbruch getrennt sind

    def compare(tokens: list[str]) -> tuple[int, Counter[str]]:
        """Gibt (fehlende Anzahl, fehlende Woerter) zurueck. Woerter, die im PDF
        ueber einen Umbruch getrennt und im Extrakt zusammengefuegt sind, gelten
        als vorhanden — das ist keine inhaltliche Abweichung."""
        budget = Counter(out)
        missing: Counter[str] = Counter()
        for i, tok in enumerate(tokens):
            if budget[tok] > 0:
                budget[tok] -= 1
                continue
            if len(tok) > 3 and tok in stream:
                res.rejoined_tokens += 1
                continue
            prev = tokens[i - 1] if i else ""
            nxt = tokens[i + 1] if i + 1 < len(tokens) else ""
            if (prev and (prev + tok) in stream) or (nxt and (tok + nxt) in stream):
                res.rejoined_tokens += 1
                continue
            missing[tok] += 1
        return sum(missing.values()), missing

    res.rejoined_tokens = 0
    total_missing, missing = compare(src_all)
    res.missing_tokens = total_missing
    res.coverage = round(100.0 * (1 - total_missing / len(src_all)), 3)
    res.missing_sample = [w for w, _ in missing.most_common(25)]

    # Seitenweise Deckung, jede Seite unabhaengig bewertet: zeigt, ob eine
    # Spalte, Tabelle oder ganze Seite im Extrakt fehlt.
    rejoined_total = res.rejoined_tokens
    for page_no, toks in pages.items():
        if not toks:
            res.per_page[page_no] = 100.0
            continue
        page_missing, _ = compare(toks)
        res.per_page[page_no] = round(100.0 * (1 - page_missing / len(toks)), 2)
    res.rejoined_tokens = rejoined_total
    res.worst_pages = sorted(res.per_page.items(), key=lambda kv: kv[1])[:5]
    return res


def office_unassigned(path: Path, md_path: Path) -> list[tuple[int, str]]:
    """Fehlende Zellen bzw. Absaetze aus Office-Quellen — Nachtrag wie beim PDF."""
    out = markdown_tokens(md_path)
    budget = Counter(out)
    stream = "".join(out)
    picked: list[tuple[int, str]] = []
    suffix = path.suffix.lower()

    def consider(page_no: int, text: str) -> None:
        toks = tokenize(text)
        if not toks:
            return
        miss = [t for t in toks if budget[t] <= 0 and not (len(t) > 3 and t in stream)]
        for t in toks:
            budget[t] -= 1
        if miss:
            picked.append((page_no, " ".join(text.split())))

    if suffix in (".xlsx", ".xlsm"):
        from openpyxl import load_workbook

        readable, _tmp = readable_office_source(path)
        wb = load_workbook(readable, read_only=True, data_only=True)
        try:
            for i, ws in enumerate(wb.worksheets, 1):
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        consider(i, " | ".join(cells))
        finally:
            wb.close()
    elif suffix == ".docx":
        from docx import Document as DocxDocument

        for para in DocxDocument(str(path)).paragraphs:
            if para.text.strip():
                consider(1, para.text)
    elif suffix == ".pptx":
        from pptx import Presentation

        for i, slide in enumerate(Presentation(str(path)).slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    consider(i, shape.text_frame.text)
    return picked


def unassigned_lines(pdf_path: Path, md_path: Path) -> list[tuple[int, str]]:
    """Zeilen des Quell-PDFs, deren Text im Extrakt fehlt — seitenweise.

    Grundlage sind die Zeilen des Docling-PDF-Backends, woertlich uebernommen.
    Damit laesst sich ein Extrakt auf 100 % Quelldeckung bringen, auch wenn das
    Tabellen- oder Layoutmodell einen Zellrest verschluckt hat.
    """
    from docling.backend.docling_parse_backend import DoclingParseDocumentBackend
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.document import InputDocument
    from docling_core.types.doc.page import TextCellUnit

    res = verify(pdf_path, md_path)
    if res.note or not res.missing_tokens:
        return []
    if pdf_path.suffix.lower() != ".pdf":
        return office_unassigned(pdf_path, md_path)

    out = markdown_tokens(md_path)
    stream = "".join(out)
    missing: Counter[str] = Counter()
    pages, _ = source_pages(pdf_path)
    budget = Counter(out)
    for toks in pages.values():
        for i, tok in enumerate(toks):
            if budget[tok] > 0:
                budget[tok] -= 1
                continue
            if len(tok) > 3 and tok in stream:
                continue
            prev = toks[i - 1] if i else ""
            nxt = toks[i + 1] if i + 1 < len(toks) else ""
            if (prev and (prev + tok) in stream) or (nxt and (tok + nxt) in stream):
                continue
            missing[tok] += 1
    if not missing:
        return []

    in_doc = InputDocument(
        path_or_stream=pdf_path, format=InputFormat.PDF,
        backend=DoclingParseDocumentBackend, filename=pdf_path.name,
    )
    backend = in_doc._backend
    picked: list[tuple[int, str]] = []
    try:
        for i in range(backend.page_count()):
            seg = backend.load_page(i).get_segmented_page()
            for cell in seg.iterate_cells(TextCellUnit.LINE):
                text = (cell.text or "").strip()
                if not text or _mask(text) == "":
                    continue
                toks = tokenize(text)
                if any(missing.get(t, 0) > 0 for t in toks):
                    for t in toks:
                        if missing.get(t, 0) > 0:
                            missing[t] -= 1
                    picked.append((i + 1, text))
    finally:
        backend.unload()
    return picked


@click.command(context_settings={"help_option_names": ["-h", "--help"]})
@click.argument("extract_md", type=click.Path(exists=True, dir_okay=False))
@click.option("--source", required=True, type=click.Path(exists=True, dir_okay=False),
              help="Quelldatei, aus der der Extrakt erzeugt wurde (PDF, XLSX, DOCX, PPTX).")
@click.option("--min-coverage", default=99.5, show_default=True,
              help="Geforderte Wortdeckung in Prozent; darunter Exit-Code 1.")
@click.option("--show-missing", default=25, show_default=True,
              help="Wie viele fehlende Woerter ausgegeben werden.")
def main(extract_md: str, source: str, min_coverage: float, show_missing: int) -> None:
    """Prueft, ob der Extrakt den Text des Quell-PDFs vollstaendig enthaelt."""
    res = verify(Path(source), Path(extract_md))
    if res.note:
        click.secho(res.note, fg="yellow")
        sys.exit(1)

    color = "green" if res.coverage >= min_coverage else "red"
    click.secho(f"Wortdeckung: {res.coverage} %  "
                f"({res.source_tokens - res.missing_tokens}/{res.source_tokens} Woerter)", fg=color)
    click.echo(f"Extrakt enthaelt {res.extract_tokens} Woerter "
               f"(Differenz kann aus Tabellen-Wiederholungen stammen).")
    if res.worst_pages:
        worst = ", ".join(f"S.{p}: {c} %" for p, c in res.worst_pages)
        click.echo(f"Schwaechste Seiten: {worst}")
    if res.boilerplate_tokens:
        click.echo(f"{res.boilerplate_tokens} Woerter aus laufenden Kopf-/Fusszeilen "
                   f"nicht gewertet (von Docling bewusst entfernt).")
    if res.rejoined_tokens:
        click.echo(f"{res.rejoined_tokens} Woerter waren im PDF ueber einen Zeilenumbruch "
                   f"getrennt und im Extrakt zusammengefuegt — als vorhanden gewertet.")
    if res.missing_tokens:
        click.echo("Fehlende Woerter (haeufigste): "
                   + ", ".join(res.missing_sample[:show_missing]))
    if res.coverage < min_coverage:
        click.secho(f"FEHLER: Deckung unter {min_coverage} % — Extrakt ist nicht vollstaendig.",
                    fg="red")
        sys.exit(1)


if __name__ == "__main__":
    main()
